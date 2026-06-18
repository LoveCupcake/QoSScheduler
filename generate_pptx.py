from pptx import Presentation
import os

prs = Presentation()

def add_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    
    title_shape.text = title
    tf = body_shape.text_frame
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

# Slide 1: Title
slide_layout = prs.slide_layouts[0] # Title slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "QoS Scheduler"
slide.shapes.placeholders[1].text = "An Advanced Network Traffic Management System for Android\nBachelor's Thesis Defense"

# Slide 2: Problem
add_slide("The Bufferbloat Dilemma in Enterprise Fleet Management", [
    "Context: Enterprises issue thousands of mobile devices to drivers (Logistics), medical staff (Healthcare), or as PoS machines.",
    "Technical Issue: Employees running background entertainment apps cause 4G bandwidth congestion (Bufferbloat), blocking critical business apps.",
    "Productivity Issue: Loss of attention during work hours due to social media usage."
])

# Slide 3: Current Limits
add_slide("Limitations of Current Market Solutions", [
    "Existing Apps: AFWall+, NetGuard, OS built-in Data Saver.",
    "Pros: Good for blocking ads and saving data quotas.",
    "Critical Flaw: Binary control only (ON/OFF firewall).",
    "Missing Features: Cannot throttle bandwidth speed or prioritize traffic queues."
])

# Slide 4: Solution
add_slide("Proposed Solution: QoS Scheduler", [
    "Operates entirely in User-Space (No Root required).",
    "Bandwidth Throttling: Precisely limits the throughput of specific applications.",
    "Weighted Fair Queuing (WFQ): Prioritizes traffic (High / Medium / Low) for mission-critical tasks.",
    "Real-time Dashboard: A centralized control panel for live telemetry and policy application."
])

# Slide 5: Arch
add_slide("Three-Plane Architecture", [
    "Data Plane: Intercepts packets via Android VpnService.",
    "Control Plane: Executes Token Bucket and WFQ algorithms.",
    "Management Plane: Web Admin dashboard built with Node.js & Chart.js."
])

# Slide 6: Tech 1
add_slide("Core Technology 1 - Data Plane", [
    "Creates a local VPN tunnel (TUN interface).",
    "Intercepts IP packets -> Parses TCP/UDP Headers.",
    "App Resolver: Maps local ports to system UIDs to identify the exact application generating the traffic (Deep Packet Inspection)."
])

# Slide 7: Tech 2
add_slide("Core Technology 2 - Control Plane", [
    "Token Bucket Algorithm: Dispenses tokens to measure and limit speed. Packets exceeding the threshold are immediately dropped.",
    "WFQ (Weighted Fair Queuing): Uses a 4:2:1 ratio. High-priority apps process tokens 4 times faster than low-priority background tasks."
])

# Slide 8: Tech 3
add_slide("Core Technology 3 - Management Plane", [
    "Lightweight Node.js backend with SQLite.",
    "1-second API polling architecture.",
    "Real-time charts (Zero animation delay) for instant visual feedback."
])

# Slide 9: Results 1
add_slide("Experimental Results: Throughput Accuracy", [
    "Token Bucket Algorithm successfully flattens traffic exactly at the Allowed BPS limit.",
    "Drop Rate (packet loss) spikes to near 100% strictly for penalized applications.",
    "Zero collateral damage: Background tasks are restricted without affecting other network traffic."
])

# Slide 10: Results 2
add_slide("Experimental Results: Bufferbloat & Latency", [
    "Weighted Fair Queuing (WFQ) successfully mitigates Bufferbloat under heavy network load.",
    "High-priority applications maintain stable, low ping (latency).",
    "Overall network fairness is preserved across the entire device."
])

# Slide 11: Demo
add_slide("Live System Demonstration", [
    "Real-time Dashboard Monitoring: Observing live traffic in the Web Admin.",
    "Applying an active Bandwidth Policy to an entertainment app.",
    "Observing the instant spike in Drop Rate and immediate throttling effect."
])

# Slide 12: Future Work
add_slide("Limitations & Future Work", [
    "Limitations: Increased battery consumption due to user-space VpnService packet processing.",
    "Future Work: Scale into a full-fledged Enterprise Fleet Management solution.",
    "Time-based Scheduling: Integrate with HR APIs to automatically throttle social media during work shifts.",
    "AI-driven QoS Automation: Integrate Machine Learning to observe app usage patterns and automatically generate dynamic bandwidth policies."
])

# Slide 13: Conclusion
slide_layout = prs.slide_layouts[0] # Title slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Thank You"
slide.shapes.placeholders[1].text = "I am now open to your questions and feedback."

output_path = os.path.join(os.getcwd(), "QoS_Scheduler_Presentation_EN_v3.pptx")
prs.save(output_path)
print(f"Generated successfully at: {output_path}")
